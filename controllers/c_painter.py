from pptx import Presentation
from pptx.util import Inches, Pt


class C_painter:
    def __init__(self, data):
        '''
        初始化C_painter
        :param data: 待绘制的数据信息
        '''
        self.data = data
        self.prs = None

    def save(self, output_path = 'undefined.pptx'):
        '''
        输出PPT
        :param output_path: 输出路径
        '''
        self.__draw()
        self.prs.save(output_path)
    
    def __draw(self):
        '''
        绘制PPT
        '''
        # 初始化
        self.prs = Presentation('templates/' + self.data.template + '.pptx')
        title_layout, contents_layout, content_layout = self.prs.slide_layouts[0], self.prs.slide_layouts[1], self.prs.slide_layouts[2]
        thanks_layout = self.prs.slide_layouts[-1]
        title_slide = self.prs.slides.add_slide(title_layout)
        contents_slide = self.prs.slides.add_slide(contents_layout)

        # 设置标题页
        title_slide.placeholders[0].text = self.data.title
        title_slide.placeholders[1].text = self.data.sub_title

        # 设置目录页
        contents = self.data.sections[0].title
        for section in self.data.sections[1:]:
            contents = contents + '\n' + section.title
        contents_slide.placeholders[11].text = contents

        # 创建内容
        for i, section in enumerate(self.data.sections):
            for j, slide in enumerate(section.slides):
                content_slide = self.prs.slides.add_slide(content_layout)
                content_slide.placeholders[0].text = f"{i+1}.{j+1} {section.title} - {slide.title}"
                content_slide.placeholders[10].text = '\n'.join(slide.texts)

        # 设置参考文献 TODO

        # 绘制结尾
        self.prs.slides.add_slide(thanks_layout)